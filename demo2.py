import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime

# ---------------------------
# CONFIGURATION
# ---------------------------
EXCEL_PATH = r"C:\Users\Divyansh\OneDrive - Casa Mia LLC\Desktop\SPec Presentation Generator\Master Product List.xlsx"
FIXED_SLIDE_IMAGES = [
    r"C:\Users\Divyansh\OneDrive - Casa Mia LLC\Desktop\SPec Presentation Generator\2nd page.jpg",
    r"C:\Users\Divyansh\OneDrive - Casa Mia LLC\Desktop\SPec Presentation Generator\3rd Page.jpg"
]
LAST_SLIDE_IMAGE = [ r"C:\Users\Divyansh\OneDrive - Casa Mia LLC\Desktop\SPec Presentation Generator\Last Page.jpg" ]
OUTPUT_PATH = r"C:\Users\Divyansh\OneDrive - Casa Mia LLC\Desktop\SPec Presentation Generator\final_presentation.pptx"
IMAGE_FOLDER = r"C:\Users\Divyansh\OneDrive - Casa Mia LLC\Desktop\SPec Presentation Generator\ERP Image Export 14-05-2025"
LOGO_PATH = r"C:\Users\Divyansh\OneDrive - Casa Mia LLC\Desktop\SPec Presentation Generator\logo.png"

TEMPLATES_PER_PRODUCT_COUNT = {
    1: 4,
    2: 3,
    3: 3,
    4: 2,
    5: 1,
    6: 1,
    7: 1,
}
today = datetime.today()
# CURRENT_NUMBER = 300   # CHANGE this number manually for each new PPT

# ---------------------------
# HELPER FUNCTIONS
# ---------------------------
def load_excel_data():
    return pd.read_excel(EXCEL_PATH)

def get_product_data(df, Code):
    return df[df["Code"].isin(Code)]

# def add_title_slide(prs, title_text, logo_path):
#     # Use Title layout (commonly 0)
#     title_layout = prs.slide_layouts[0]
#     slide = prs.slides.add_slide(title_layout)
#     slide.shapes.title.text = title_text

#     # Add BIG logo at bottom right
#     LOGO_WIDTH = Inches(3.5)
#     LOGO_HEIGHT = Inches(1.5)
#     margin_right = Inches(0.5)
#     margin_bottom = Inches(0.3)
#     SLIDE_WIDTH = prs.slide_width
#     SLIDE_HEIGHT = prs.slide_height

#     left = SLIDE_WIDTH - LOGO_WIDTH - margin_right
#     top = SLIDE_HEIGHT - LOGO_HEIGHT - margin_bottom

#     if os.path.exists(logo_path):
#         slide.shapes.add_picture(logo_path, left, top, width=LOGO_WIDTH, height=LOGO_HEIGHT)
#     else:
#         print("⚠️ Logo image not found for title slide!")

def get_next_counter(counter_file, start_from):
    if os.path.exists(counter_file):
        with open(counter_file, 'r') as f:
            try:
                num = int(f.read().strip())
            except:
                num = start_from
    else:
        num = start_from
    return num

def increment_counter(counter_file, num):
    with open(counter_file, 'w') as f:
        f.write(str(num + 1))


def add_title_slide(
    prs, 
    main_title="PROPOSED SELECTIONS", 
    subtitle="SANITARYWARE |", 
    date=" APRIL 15,2025",
    project="MAF - MOE", 
    logo_path=None
):
    # Use a blank slide for full control
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[-1]

    slide = prs.slides.add_slide(blank_layout)
    SLIDE_WIDTH = prs.slide_width
    SLIDE_HEIGHT = prs.slide_height

    # --- Main Title (big, red, right) ---
    title_box = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(11),  # x from left
        Inches(0.8),                 # y from top
        Inches(10),                  # width
        Inches(3.5)                  # height
    )
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = main_title
    p.font.size = Pt(100)
    p.font.name = 'Belleza'           # or your preferred font
    p.font.color.rgb = RGBColor(225, 25, 25)  # Red
    p.alignment = PP_ALIGN.RIGHT

    # --- Subtitle (smaller, red, right, below title) ---
    subtitle_box = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(11),
        Inches(4.5),
        Inches(10),
        Inches(0.4)
    )
    tf2 = subtitle_box.text_frame
    tf2.clear()
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(30)
    p2.font.name = 'Poppins Light'
    p2.font.color.rgb = RGBColor(225, 25, 25)  # Red
    p2.alignment = PP_ALIGN.RIGHT

    # --- Subtitle (smaller, red, right, below title) ---
    date_box = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(11),
        Inches(5),
        Inches(10),
        Inches(0.5)
    )
    tf2 = date_box.text_frame
    tf2.clear()
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = date
    p2.font.size = Pt(30)
    p2.font.name = 'Poppins Light'
    p2.font.color.rgb = RGBColor(225, 25, 25)  # Red
    p2.alignment = PP_ALIGN.RIGHT

    # --- Project (orange, right, lower on slide) ---
    proj_box = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(7),    # right side
        Inches(8.5),                # a bit lower down
        Inches(6),
        Inches(1)
    )
    tf3 = proj_box.text_frame
    tf3.clear()
    p3 = tf3.paragraphs[0]
    p3.text = project
    p3.font.size = Pt(50)
    p3.font.name = 'Belleza'
    p3.font.color.rgb = RGBColor(237, 137, 39)  # Orange
    p3.alignment = PP_ALIGN.RIGHT

    # --- Logo (bottom right) ---
    LOGO_WIDTH = Inches(5)
    LOGO_HEIGHT = Inches(2)
    margin_right = Inches(0.3)
    margin_bottom = Inches(0.1)
    left = SLIDE_WIDTH - LOGO_WIDTH - margin_right
    top = SLIDE_HEIGHT - LOGO_HEIGHT - margin_bottom
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, left, top, width=LOGO_WIDTH, height=LOGO_HEIGHT)
    else:
        print("⚠️ Logo image not found for title slide!")

def add_fixed_slide_from_image(prs, image_path):
    # Use a safe blank layout (usually the last one)
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide.shapes.add_picture(image_path, 0, 0, width=slide_width, height=slide_height)

from layout_manager import get_layout_positions

def add_product_slide(prs, products, title, template_number):
    # Use a safe blank layout
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # Add page title manually at the top
    # title_box = slide.shapes.add_textbox(Inches(3.335), Inches(0.3), Inches(20), Inches(1)) #LEFT TOP WIDDTH HEIGHT
    slide.shapes.add_picture( r"C:\Users\Divyansh\OneDrive - Casa Mia LLC\Desktop\SPec Presentation Generator\Heading Background.png", Inches(9.835), Inches(0), Inches(7), Inches(1.25))
    title_box = slide.shapes.add_textbox(Inches(9.835), Inches(0.25), Inches(7), Inches(1.25))

    tf_title = title_box.text_frame
    tf_title.clear()
    p = tf_title.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.name = 'Poppins Light'
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    product_count = len(products)
    positions, image_size, text_height = get_layout_positions(product_count, template_number)

    for i, (_, product) in enumerate(products.iterrows()):
        if i >= len(positions):
            break

        left, top = positions[i]
        img_filename = f"{str(product['Code'])}.jpeg"  # or .png if you use png
        img_path = os.path.join(IMAGE_FOLDER, img_filename)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, width=image_size, height=image_size)
        else:
            print(f"⚠️ Image not found for ERP Code: {product['Code']}")

        # Add textbox under image
        textbox = slide.shapes.add_textbox(left, top + image_size + Inches(0.3), image_size + Inches(0.5), text_height)
        tf = textbox.text_frame
        tf.clear()
        tf.word_wrap = True
        # text = f"{product['Code']} - {product['Description']}\n"
        # for line in text.strip().split('\n'):
        #     p = tf.add_paragraph()
        #     p.text = line
        #     p.font.size = Pt(18)
        #     p.font.name = 'Poppins Light'
        #     p.font.color.rgb = RGBColor(0, 0, 0)
        #     p.alignment = PP_ALIGN.JUSTIFY
        
        code = product['Code']
        description = product['Description']

        # First line: Code (bold) + Description (normal)
        p = tf.add_paragraph()
        run_code = p.add_run()
        run_code.text = f"{code} - "
        run_code.font.bold = True  # Make the code bold
        run_code.font.size = Pt(18)
        run_code.font.name = 'Poppins Light'
        run_code.font.color.rgb = RGBColor(0, 0, 0)

        run_desc = p.add_run()
        run_desc.text = description
        run_desc.font.bold = False
        run_desc.font.size = Pt(18)
        run_desc.font.name = 'Poppins Light'
        run_desc.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.JUSTIFY

    # --- Add logo at bottom right ---
    LOGO_WIDTH = Inches(2.5)
    LOGO_HEIGHT = Inches(1.2)
    SLIDE_WIDTH = prs.slide_width
    SLIDE_HEIGHT = prs.slide_height
    margin_right = Inches(0.3)
    margin_bottom = Inches(0.1)
    left = SLIDE_WIDTH - LOGO_WIDTH - margin_right
    top = SLIDE_HEIGHT - LOGO_HEIGHT - margin_bottom
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, left, top, width=LOGO_WIDTH, height=LOGO_HEIGHT)
    else:
        print("⚠️ Logo image not found for product slide!")

# ---------------------------
# MAIN EXECUTION
# ---------------------------
def main():
    df = load_excel_data()
    prs = Presentation()  # Start blank
    prs.slide_width = Inches(26.67)
    prs.slide_height = Inches(15)

    # Step 1: Add title slide
    # title = input("Enter the title of the presentation: ")
    # add_title_slide(prs, title, LOGO_PATH)
    # subtitletext = input("Enter the subtitle : ") 
    formatted_date = today.strftime("%B %d, %Y")
    projecttext = input("Enter the project name: ")

    add_title_slide(
    prs,
    main_title="PROPOSED SELECTIONS",
    # subtitle= subtitletext + "  |",
    subtitle= "SANITARYWARE |",
    date= formatted_date,
    project= projecttext,
    logo_path=LOGO_PATH
)



    # Step 2: Insert fixed slides as images
    for img_path in FIXED_SLIDE_IMAGES:
        add_fixed_slide_from_image(prs, img_path)

    # Step 3: Add product slides interactively
    while True:
        page_title = input("\nEnter the title for this product page: ")
        product_count = int(input("How many products on this slide? (1/2/3/4/5/6/7): "))

        template_count = TEMPLATES_PER_PRODUCT_COUNT.get(product_count)
        if template_count is None:
            print(f"⚠️ No templates defined for product count {product_count}.")
            continue

        if template_count == 1:
            print("Only 1 template available for this product count. Selecting Template 1.")
            template_number = 1
        else:
            template_options = "/".join(str(i+1) for i in range(template_count))
            template_number = int(input(f"Choose layout template ({template_options}): "))
            if not (1 <= template_number <= template_count):
                print("Invalid template selected.")
                continue

        Code = []
        for i in range(product_count):
            Code.append(input(f"Enter Code #{i+1}: ").strip())

        selected_products = get_product_data(df, Code)
        if selected_products.empty:
            print("⚠️ No matching ERP IDs found. Try again.")
            continue

        add_product_slide(prs, selected_products, page_title, template_number)

        cont = input("Do you want to add another product page? (yes/no): ").lower()
        if cont != "yes":
            break
    
    # Step 2: Insert fixed slides as images
    for img_path in LAST_SLIDE_IMAGE:
        add_fixed_slide_from_image(prs, img_path)
    
    # Step 4: Save presentation

    year_str = f"YR{today.strftime('%y')}"
    prefix = "PP"
    COUNTER_FILE = os.path.join(os.path.dirname(OUTPUT_PATH), "ppt_counter.txt")
    START_FROM = 300  # Change this if you want a new base

    current_number = get_next_counter(COUNTER_FILE, START_FROM)
    file_number_str = f"{current_number:04d}"  # zero-padded

    output_path = os.path.join(
        os.path.dirname(OUTPUT_PATH),
        f"{prefix}-{year_str}-{file_number_str}.pptx"
    )
    prs.save(output_path)
    # print(f"\n✅ Presentation saved at: {output_path}")
    # print("⚠️ If any image was missing, it has been reported above.")
    ppt_filename = os.path.basename(output_path)  # Extracts just the filename

    print(f"\n✅ Presentation saved at: {output_path}")
    print(f"   📁 File name: {ppt_filename}")
    print("⚠️ If any image was missing, it has been reported above.")
    increment_counter(COUNTER_FILE, current_number)


    # safe_title = "".join(c for c in title if c.isalnum() or c in " _-").strip()
    # output_path = os.path.join(
    #     os.path.dirname(OUTPUT_PATH), f"{safe_title} - final presentation.pptx"
    # )
    # prs.save(output_path)
    # print(f"\n✅ Presentation saved at: {output_path}")
    # print("⚠️ If any image was missing, it has been reported above.")

if __name__ == "__main__":
    main()
